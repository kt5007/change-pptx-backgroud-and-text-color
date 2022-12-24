import pptx
import collections.abc
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

if __name__ == '__main__':
    prs = pptx.Presentation('.pptx')
    print("1ページ目")
    slide = prs.slides[2]
    k=0
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255,255,255)
    for shape in slide.shapes:
        # print(shape.name)
        print(k)
        k+=1
        if not shape.has_text_frame: 
            continue
        textFrame = shape.text_frame                   # Get TextFrame from various Shape objects
        textFrame.text = shape.text_frame.text      # Set text in TextFrame
        for paragraph in textFrame.paragraphs:
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE    # paragraph text size auto-adjustment
            paragraph.word_wrap = True                              # Enable auto text wrapping
    
    prs.save('Art2_List1_result.pptx')
